from pptx import Presentation
from gtts import gTTS
from playsound import playsound
import os

def extract_text_from_powerpoint(pptx_file_path):
    prs = Presentation(pptx_file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def text_to_speech(text, lang='vi'):
    tts = gTTS(text, lang=lang)
    tts.save('output.mp3')

# Đường dẫn tới tệp PowerPoint cần trích xuất văn bản
pptx_file_path = "D:\\Python\\Nhóm 2.pptx"

# Trích xuất văn bản từ tệp PowerPoint
extracted_text = extract_text_from_powerpoint(pptx_file_path)

# Chuyển đổi văn bản thành giọng nói và lưu thành tệp âm thanh
text_to_speech(extracted_text)

# Phát âm thanh đã tạo
playsound('output.mp3')